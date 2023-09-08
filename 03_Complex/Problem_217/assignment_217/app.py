from flask import Flask, render_template, request, send_file
from pptx import Presentation
import io

app = Flask(__name__)

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        # Get data from the form
        title = request.form.get('title')
        content = request.form.get('content')

        # Create a PowerPoint presentation
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = title
        tf = body_shape.text_frame
        tf.text = content

        # Save the presentation to a byte stream
        pptx_buffer = io.BytesIO()
        presentation.save(pptx_buffer)
        pptx_buffer.seek(0)

        # Return the PowerPoint file for download
        return send_file(
            pptx_buffer,
            attachment_filename='generated_presentation.pptx',
            as_attachment=True
        )

    return render_template('index.html')

if __name__ == '__main__':
    app.run(debug=True)
